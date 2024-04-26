# !/usr/bin/env python3
# -*- coding:utf-8 -*-

# @Time    : 2024/3/18 17:34
# @Author  : wangchongshi
# @Email   : wangchongshi.wcs@antgroup.com
# @FileName: pptx_parser.py
from pathlib import Path
from typing import Dict, List, Optional

from agentuniverse.agent.action.knowledge.reader.reader import Reader
from agentuniverse.agent.action.knowledge.store.document import Document


class PptxReader(Reader):
    """Pptx reader."""

    def load_data(self, file: Path, ext_info: Optional[Dict] = None) -> List[Document]:
        """Parse the pptx file.

        Note:
            `python-pptx` is required to read PPTX files: `pip install python-pptx`
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required to read pptx files: `pip install python-pptx`"
            )
        presentation = Presentation(file)
        document_list = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    metadata = {"slide_number": slide_number, "file_name": file.name}
                    if ext_info is not None:
                        metadata.update(ext_info)
                    # Extract the text from the shape
                    document_list.append(Document(text=shape.text, metadata=metadata))
        return document_list
